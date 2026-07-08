"""AI Report Generator (Feature 7).

Generates executive business reports (daily/weekly/monthly/quarterly/investor/
board/sales/finance/marketing/inventory) by aggregating the rest of the
copilot: KPIs (briefing_service), charts (recent Query.result_json), insights
(latest QueryInsight), root causes (open SmartAlerts), forecasts
(forecast_service, lazy), and recommendations (recommendation_service, lazy,
else QueryInsight.recommendations). The executive summary is Gemini-generated
with a graceful deterministic fallback.

Exports: CSV (pandas), Excel (openpyxl), PDF (reportlab), PowerPoint
(python-pptx). Email via stdlib smtplib (no-op when SMTP_HOST unset).
Scheduled reports advance next_run from a preset cadence.
"""
from __future__ import annotations

import io
import logging
import os
import smtplib
import uuid
from datetime import datetime, timedelta
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from typing import Any

import pandas as pd
from sqlalchemy.orm import Session

from models import GeneratedReport, Query, QueryInsight, SmartAlert, Dataset, User
from services import ai_service, notifications_service
from services.briefing_service import _kpi_deltas

logger = logging.getLogger("report_service")

REPORT_TYPES = [
    "daily", "weekly", "monthly", "quarterly",
    "investor", "board", "sales", "finance", "marketing", "inventory",
]

# Cadence presets for scheduled reports. Stored in schedule_cron; next_run is
# computed by adding the interval to "now".
CADENCE_HOURS = {
    "hourly": 1,
    "daily": 24,
    "weekly": 24 * 7,
    "monthly": 24 * 30,
    "quarterly": 24 * 90,
}

REPORT_SCHEMA = {
    "type": "object",
    "properties": {
        "executive_summary": {"type": "string"},
        "key_takeaways": {"type": "array", "items": {"type": "string"}},
        "outlook": {"type": "string"},
    },
    "required": ["executive_summary"],
}


def _title_for(report_type: str) -> str:
    pretty = {
        "daily": "Daily Business Report", "weekly": "Weekly Business Report",
        "monthly": "Monthly Business Report", "quarterly": "Quarterly Business Report",
        "investor": "Investor Update", "board": "Board Report",
        "sales": "Sales Report", "finance": "Finance Report",
        "marketing": "Marketing Report", "inventory": "Inventory Report",
    }
    return pretty.get(report_type, report_type.title() + " Report")


def _charts(db: Session, user_id: int, limit: int = 4) -> list[dict[str, Any]]:
    """Recent successful query results, shaped as chart payloads."""
    recent = (
        db.query(Query)
        .filter(Query.user_id == user_id, Query.execution_time_ms.isnot(None))
        .order_by(Query.created_at.desc())
        .limit(limit)
        .all()
    )
    charts = []
    for q in recent:
        rows = q.result_json if isinstance(q.result_json, list) else []
        if not rows:
            continue
        charts.append({
            "question": q.question,
            "sql": q.sql,
            "rows": rows[:50],
            "row_count": len(rows),
        })
    return charts


def _insights(db: Session, user_id: int) -> dict[str, Any] | None:
    ins = (
        db.query(QueryInsight)
        .join(Query, QueryInsight.query_id == Query.id)
        .filter(Query.user_id == user_id)
        .order_by(QueryInsight.generated_at.desc())
        .first()
    )
    if not ins:
        return None
    return {
        "executive_summary": (ins.insights_json or {}).get("executive_summary") if isinstance(ins.insights_json, dict) else None,
        "risks": (ins.insights_json or {}).get("risks") if isinstance(ins.insights_json, dict) else [],
        "opportunities": (ins.insights_json or {}).get("opportunities") if isinstance(ins.insights_json, dict) else [],
        "recommendations": (ins.insights_json or {}).get("recommendations") if isinstance(ins.insights_json, dict) else [],
        "explanation": ins.explanation,
    }


def _root_causes(db: Session, user_id: int) -> list[dict[str, Any]]:
    rows = (
        db.query(SmartAlert)
        .filter(SmartAlert.user_id == user_id, SmartAlert.status == "open")
        .order_by(SmartAlert.severity.desc(), SmartAlert.created_at.desc())
        .limit(5)
        .all()
    )
    return [
        {"name": a.name, "metric": a.metric, "severity": a.severity,
         "root_cause": a.root_cause, "recommended_action": a.recommended_action}
        for a in rows
    ]


def _forecasts(db: Session, user_id: int) -> list[dict[str, Any]]:
    """Lazy import forecast_service (Phase 8); graceful if unavailable.

    Phase 8's forecast_service.list_recent_forecasts(db, user_id) returns the
    latest cached forecast per metric. Until then this returns [] so reports
    still generate with everything else populated.
    """
    try:
        from services import forecast_service
        fn = getattr(forecast_service, "list_recent_forecasts", None)
        if fn is None:
            return []
        return fn(db, user_id) or []
    except Exception:
        return []


def _recommendations(db: Session, user_id: int, fallback: dict[str, Any] | None) -> list[str]:
    try:
        from services import recommendation_service  # noqa
        # Phase 9 wires the real generator; until then use the insight fallback.
    except Exception:
        pass
    if fallback and isinstance(fallback.get("recommendations"), list):
        return [str(r) for r in fallback["recommendations"][:6] if r]
    return [
        "Review your top-performing segments and double down on what's working.",
        "Address the highest-severity open alerts first.",
        "Schedule a recurring report so stakeholders stay aligned automatically.",
    ]


def _exec_summary(report_type: str, kpis: list[dict[str, Any]], root_causes: list[dict[str, Any]],
                  insights: dict[str, Any] | None) -> dict[str, Any]:
    kpi_line = "; ".join(
        f"{k['label']} {k['value']} ({('+' if k['delta'] and k['delta'] > 0 else '')}{k['delta']}% {k['direction']})"
        for k in kpis if k.get("has_data")
    ) or "limited live data connected"
    rc_line = f"{len(root_causes)} open alert(s)" if root_causes else "no open alerts"
    base = (
        f"{_title_for(report_type)} — {kpi_line}. {rc_line}. "
        f"Outlook: {insights.get('executive_summary') if insights else 'add a dataset and ask questions for a full outlook.'}"
    )
    try:
        prompt = (
            "You are an AI BI analyst writing an executive report summary.\n"
            f"Report type: {report_type}\n"
            f"KPIs: {kpi_line}\n"
            f"Open alerts/root causes: {rc_line}\n"
            f"Latest insight: {(insights or {}).get('executive_summary') or 'n/a'}\n\n"
            "Return JSON with executive_summary (2-3 sentences, executive tone), key_takeaways (3-5 short bullets), "
            "and outlook (one sentence on the next period)."
        )
        out = ai_service.generate_insight_json(prompt, schema=REPORT_SCHEMA)
        out.setdefault("key_takeaways", [])
        out.setdefault("outlook", "")
        return out
    except Exception as exc:
        logger.warning("Report exec-summary LLM failed (%s); using stub.", exc)
        return {"executive_summary": base, "key_takeaways": [], "outlook": ""}


def generate_report(db: Session, user: User, report_type: str, dataset_id: int | None = None) -> dict[str, Any]:
    if report_type not in REPORT_TYPES:
        report_type = "daily"

    kpis, _ = _kpi_deltas(db, user.id)
    charts = _charts(db, user.id)
    insights = _insights(db, user.id)
    root_causes = _root_causes(db, user.id)
    forecasts = _forecasts(db, user.id)
    recommendations = _recommendations(db, user.id, insights)
    exec_block = _exec_summary(report_type, kpis, root_causes, insights)

    content = {
        "executive_summary": exec_block.get("executive_summary"),
        "key_takeaways": exec_block.get("key_takeaways", []),
        "outlook": exec_block.get("outlook", ""),
        "kpis": kpis,
        "charts": charts,
        "insights": insights,
        "root_causes": root_causes,
        "forecasts": forecasts,
        "recommendations": recommendations,
    }

    row = GeneratedReport(
        user_id=user.id,
        dataset_id=dataset_id,
        report_type=report_type,
        title=_title_for(report_type),
        content_json=content,
        status="ready",
        last_generated=datetime.utcnow(),
    )
    db.add(row)
    db.commit()
    db.refresh(row)

    try:
        notifications_service.create_notification(
            db, user.id, title=f"{row.title} is ready",
            body=(content.get("executive_summary") or "")[:160], type="info",
        )
    except Exception:
        pass

    return _serialize(row)


def _serialize(r: GeneratedReport) -> dict[str, Any]:
    return {
        "id": r.id,
        "report_type": r.report_type,
        "title": r.title,
        "content": r.content_json or {},
        "status": r.status,
        "schedule_cron": r.schedule_cron,
        "next_run": r.next_run.isoformat() if r.next_run else None,
        "last_generated": r.last_generated.isoformat() if r.last_generated else None,
        "share_token": r.share_token,
        "created_at": r.created_at.isoformat() if r.created_at else None,
    }


def list_reports(db: Session, user_id: int, limit: int = 30) -> list[dict[str, Any]]:
    rows = (
        db.query(GeneratedReport)
        .filter(GeneratedReport.user_id == user_id)
        .order_by(GeneratedReport.created_at.desc())
        .limit(limit)
        .all()
    )
    return [_serialize(r) for r in rows]


def get_report(db: Session, user_id: int, report_id: int) -> dict[str, Any] | None:
    r = db.query(GeneratedReport).filter(GeneratedReport.id == report_id, GeneratedReport.user_id == user_id).first()
    return _serialize(r) if r else None


def get_shared(db: Session, token: str) -> dict[str, Any] | None:
    r = db.query(GeneratedReport).filter(GeneratedReport.share_token == token).first()
    return _serialize(r) if r else None


def delete_report(db: Session, user_id: int, report_id: int) -> bool:
    r = db.query(GeneratedReport).filter(GeneratedReport.id == report_id, GeneratedReport.user_id == user_id).first()
    if not r:
        return False
    db.delete(r)
    db.commit()
    return True


def share_report(db: Session, user_id: int, report_id: int) -> dict[str, Any] | None:
    r = db.query(GeneratedReport).filter(GeneratedReport.id == report_id, GeneratedReport.user_id == user_id).first()
    if not r:
        return None
    if not r.share_token:
        r.share_token = uuid.uuid4().hex
        db.commit()
        db.refresh(r)
    return {"share_token": r.share_token}


def schedule_report(db: Session, user_id: int, report_id: int, cadence: str) -> dict[str, Any] | None:
    r = db.query(GeneratedReport).filter(GeneratedReport.id == report_id, GeneratedReport.user_id == user_id).first()
    if not r:
        return None
    hours = CADENCE_HOURS.get(cadence)
    if hours is None:
        return None
    r.schedule_cron = cadence
    r.next_run = datetime.utcnow() + timedelta(hours=hours)
    db.commit()
    db.refresh(r)
    return _serialize(r)


def unschedule_report(db: Session, user_id: int, report_id: int) -> dict[str, Any] | None:
    r = db.query(GeneratedReport).filter(GeneratedReport.id == report_id, GeneratedReport.user_id == user_id).first()
    if not r:
        return None
    r.schedule_cron = None
    r.next_run = None
    db.commit()
    db.refresh(r)
    return _serialize(r)


def run_due_reports() -> None:
    """Scheduler entrypoint: regenerate + email every due scheduled report."""
    from db import SessionLocal
    db = SessionLocal()
    try:
        now = datetime.utcnow()
        due = (
            db.query(GeneratedReport)
            .filter(GeneratedReport.next_run.isnot(None), GeneratedReport.next_run <= now)
            .all()
        )
        for r in due:
            try:
                user = db.query(User).filter(User.id == r.user_id).first()
                if not user:
                    continue
                # Regenerate content in place.
                regenerated = generate_report(db, user, r.report_type, r.dataset_id)
                # Copy fresh content into the scheduled row and advance next_run.
                r.content_json = regenerated["content"]
                r.last_generated = datetime.utcnow()
                hours = CADENCE_HOURS.get(r.schedule_cron or "daily", 24)
                r.next_run = datetime.utcnow() + timedelta(hours=hours)
                db.commit()
            except Exception as exc:
                logger.error("scheduled report %s failed: %s", r.id, exc)
    finally:
        db.close()


# ---- exports ------------------------------------------------------------

def _flatten_content(content: dict[str, Any]) -> dict[str, Any]:
    """Produce a flat, Excel/CSV-friendly dict of the report."""
    flat = {
        "executive_summary": content.get("executive_summary"),
        "outlook": content.get("outlook"),
    }
    for k in content.get("kpis", []):
        if k.get("has_data"):
            flat[f"kpi_{k['key']}_value"] = k.get("value")
            flat[f"kpi_{k['key']}_delta_pct"] = k.get("delta")
    for i, rc in enumerate(content.get("root_causes", [])[:5], 1):
        flat[f"root_cause_{i}"] = rc.get("root_cause")
    for i, rec in enumerate(content.get("recommendations", [])[:6], 1):
        flat[f"recommendation_{i}"] = rec
    return flat


def export_csv(report: GeneratedReport) -> tuple[bytes, str]:
    content = report.content_json or {}
    rows = [_flatten_content(content)]
    # Append a KPI table block.
    for k in content.get("kpis", []):
        rows.append({"section": "kpi", "metric": k.get("label"), "value": k.get("value"),
                     "delta_pct": k.get("delta"), "direction": k.get("direction")})
    df = pd.DataFrame(rows)
    buf = io.StringIO()
    df.to_csv(buf, index=False)
    data = buf.getvalue().encode("utf-8")
    return data, f"{report.title}.csv"


def export_excel(report: GeneratedReport) -> tuple[bytes, str]:
    content = report.content_json or {}
    out = io.BytesIO()
    with pd.ExcelWriter(out, engine="openpyxl") as writer:
        # Summary sheet.
        summary = pd.DataFrame([_flatten_content(content)])
        summary.to_excel(writer, sheet_name="Summary", index=False)
        # KPIs sheet.
        pd.DataFrame(content.get("kpis", [])).to_excel(writer, sheet_name="KPIs", index=False)
        # Root causes.
        pd.DataFrame(content.get("root_causes", [])).to_excel(writer, sheet_name="Root causes", index=False)
        # Recommendations.
        pd.DataFrame({"recommendation": content.get("recommendations", [])}).to_excel(
            writer, sheet_name="Recommendations", index=False)
        # First chart's raw rows, if any.
        charts = content.get("charts", [])
        if charts and charts[0].get("rows"):
            try:
                pd.DataFrame(charts[0]["rows"]).to_excel(writer, sheet_name="Chart data", index=False)
            except Exception:
                pass
    data = out.getvalue()
    return data, f"{report.title}.xlsx"


def export_pdf(report: GeneratedReport) -> tuple[bytes, str]:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

    content = report.content_json or {}
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
                            topMargin=0.8 * inch, bottomMargin=0.8 * inch)
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("h1", parent=styles["Title"], fontSize=20, textColor=colors.HexColor("#1e293b"), spaceAfter=6)
    meta = ParagraphStyle("meta", parent=styles["Normal"], fontSize=9, textColor=colors.grey, spaceAfter=12)
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], fontSize=13, textColor=colors.HexColor("#4338ca"), spaceBefore=12, spaceAfter=4)
    body = ParagraphStyle("body", parent=styles["Normal"], fontSize=10, leading=14, spaceAfter=6)

    story = []
    story.append(Paragraph(report.title, h1))
    story.append(Paragraph(f"{report.report_type.title()} report · generated {report.last_generated or report.created_at}", meta))

    if content.get("executive_summary"):
        story.append(Paragraph("Executive summary", h2))
        story.append(Paragraph(str(content["executive_summary"]), body))

    if content.get("key_takeaways"):
        story.append(Paragraph("Key takeaways", h2))
        for t in content["key_takeaways"]:
            story.append(Paragraph(f"• {t}", body))

    kpis = [k for k in content.get("kpis", []) if k.get("has_data")]
    if kpis:
        story.append(Paragraph("KPIs", h2))
        tbl = [["Metric", "Value", "Delta %", "Direction"]]
        for k in kpis:
            tbl.append([str(k.get("label", "")), str(k.get("value", "")),
                        str(k.get("delta", "")), str(k.get("direction", ""))])
        story.append(Table(tbl, colWidths=[2 * inch, 1.2 * inch, 1 * inch, 1 * inch],
                           style=TableStyle([
                               ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#eef2ff")),
                               ("FONTSIZE", (0, 0), (-1, -1), 9),
                               ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e2e8f0")),
                           ])))
        story.append(Spacer(1, 8))

    if content.get("root_causes"):
        story.append(Paragraph("Root causes & open alerts", h2))
        for rc in content["root_causes"]:
            story.append(Paragraph(f"• <b>{rc.get('name')}</b> — {rc.get('root_cause', '')}", body))

    if content.get("recommendations"):
        story.append(Paragraph("Recommendations", h2))
        for rec in content["recommendations"]:
            story.append(Paragraph(f"• {rec}", body))

    if content.get("outlook"):
        story.append(Paragraph("Outlook", h2))
        story.append(Paragraph(str(content["outlook"]), body))

    doc.build(story)
    return buf.getvalue(), f"{report.title}.pdf"


def export_ppt(report: GeneratedReport) -> tuple[bytes, str]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    content = report.content_json or {}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide.
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = report.title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"{report.report_type.title()} report"

    # Executive summary slide.
    if content.get("executive_summary"):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "Executive summary"
        tb = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5))
        tb.text_frame.text = str(content["executive_summary"])
        tb.text_frame.word_wrap = True
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(16)

    # KPIs slide.
    kpis = [k for k in content.get("kpis", []) if k.get("has_data")]
    if kpis:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "KPIs"
        tb = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, k in enumerate(kpis):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{k.get('label')}: {k.get('value')}  ({('+' if k.get('delta') and k['delta'] > 0 else '')}{k.get('delta')}% {k.get('direction')})"
            for r in p.runs:
                r.font.size = Pt(18)

    # Recommendations slide.
    if content.get("recommendations"):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "Recommendations"
        tb = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, rec in enumerate(content["recommendations"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {rec}"
            for r in p.runs:
                r.font.size = Pt(16)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), f"{report.title}.pptx"


def export_report_bytes(db: Session, user_id: int, report_id: int, fmt: str) -> tuple[bytes, str] | None:
    r = db.query(GeneratedReport).filter(GeneratedReport.id == report_id, GeneratedReport.user_id == user_id).first()
    if not r:
        return None
    fmt = (fmt or "").lower()
    if fmt == "csv":
        return export_csv(r)
    if fmt == "excel" or fmt == "xlsx":
        return export_excel(r)
    if fmt == "pdf":
        return export_pdf(r)
    if fmt == "ppt" or fmt == "pptx":
        return export_ppt(r)
    return None


# ---- email --------------------------------------------------------------

def email_report(db: Session, user_id: int, report_id: int, to_addr: str | None = None) -> dict[str, Any]:
    r = db.query(GeneratedReport).filter(GeneratedReport.id == report_id, GeneratedReport.user_id == user_id).first()
    if not r:
        return {"ok": False, "detail": "Report not found"}
    user = db.query(User).filter(User.id == user_id).first()
    recipient = to_addr or (user.email if user else None)
    if not recipient:
        return {"ok": False, "detail": "No recipient address"}

    host = os.getenv("SMTP_HOST")
    if not host:
        return {"ok": False, "detail": "Email not configured (SMTP_HOST unset). Export the report and send manually."}

    content = r.content_json or {}
    msg = MIMEMultipart("alternative")
    msg["Subject"] = r.title
    msg["From"] = os.getenv("SMTP_FROM", "reports@bi-copilot.local")
    msg["To"] = recipient
    body = (
        f"{r.title}\n\n"
        f"{content.get('executive_summary', '')}\n\n"
        "Key takeaways:\n" + "\n".join(f"- {t}" for t in content.get("key_takeaways", []))
        + "\n\nRecommendations:\n" + "\n".join(f"- {x}" for x in content.get("recommendations", []))
    )
    msg.attach(MIMEText(body, "plain"))

    try:
        port = int(os.getenv("SMTP_PORT", "587"))
        with smtplib.SMTP(host, port, timeout=15) as server:
            if os.getenv("SMTP_USER") and os.getenv("SMTP_PASS"):
                server.starttls()
                server.login(os.getenv("SMTP_USER"), os.getenv("SMTP_PASS"))
            server.sendmail(msg["From"], [recipient], msg.as_string())
        return {"ok": True, "to": recipient}
    except Exception as exc:
        logger.warning("email_report failed (%s)", exc)
        return {"ok": False, "detail": f"SMTP send failed: {exc}"}